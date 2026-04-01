import os
from dataclasses import dataclass, field
from app.utils.file_utils import get_file_extension


@dataclass
class DocumentSection:
    text: str
    metadata: dict = field(default_factory=dict)


def process_file(file_path: str, filename: str) -> list[DocumentSection]:
    ext = get_file_extension(filename)
    processors = {
        ".pdf": _process_pdf,
        ".docx": _process_docx,
        ".xlsx": _process_xlsx,
        ".pptx": _process_pptx,
        ".txt": _process_txt,
        ".md": _process_txt,
    }
    processor = processors.get(ext)
    if processor is None:
        raise ValueError(f"Unsupported file type: {ext}")
    return processor(file_path, filename)


def _process_pdf(file_path: str, filename: str) -> list[DocumentSection]:
    import pdfplumber

    sections = []
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text and text.strip():
                sections.append(DocumentSection(
                    text=text.strip(),
                    metadata={"source": filename, "page": i + 1},
                ))
    return sections


def _process_docx(file_path: str, filename: str) -> list[DocumentSection]:
    from docx import Document

    doc = Document(file_path)
    parts = []

    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())

    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip(" |"):
                parts.append(row_text)

    full_text = "\n".join(parts)
    if full_text.strip():
        return [DocumentSection(text=full_text, metadata={"source": filename})]
    return []


def _process_xlsx(file_path: str, filename: str) -> list[DocumentSection]:
    from openpyxl import load_workbook

    wb = load_workbook(file_path, read_only=True, data_only=True)
    sections = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
            if row_text.strip(" |"):
                rows.append(row_text)
        if rows:
            text = f"Sheet: {sheet_name}\n" + "\n".join(rows)
            sections.append(DocumentSection(
                text=text,
                metadata={"source": filename, "sheet": sheet_name},
            ))

    wb.close()
    return sections


def _process_pptx(file_path: str, filename: str) -> list[DocumentSection]:
    from pptx import Presentation

    prs = Presentation(file_path)
    sections = []

    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text.strip())
        if texts:
            sections.append(DocumentSection(
                text="\n".join(texts),
                metadata={"source": filename, "slide": i + 1},
            ))

    return sections


def _process_txt(file_path: str, filename: str) -> list[DocumentSection]:
    text = None
    for encoding in ["utf-8", "gbk", "latin-1"]:
        try:
            with open(file_path, "r", encoding=encoding) as f:
                text = f.read()
            break
        except (UnicodeDecodeError, ValueError):
            continue

    if text and text.strip():
        return [DocumentSection(text=text.strip(), metadata={"source": filename})]
    return []
